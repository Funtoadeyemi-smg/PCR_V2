from __future__ import annotations

import os
import re
from typing import Dict, Iterable, Optional, Sequence

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER


class PowerPointProcessor:
    """Handles PowerPoint modifications for the consolidated template."""

    _PLACEHOLDER_PATTERN = re.compile(r"\{[^{}]+\}")
    _CHANNEL_DISPLAY = {
        "meta": "meta",
        "pin": "pinterest",
        "tik": "tiktok",
    }
    _IMAGE_LAYOUT_OVERRIDES = {
        "{table_of_contents_picture}": {
            "replace_existing_picture": True,
            "width_factor": 0.72,
            "height_factor": 0.55,
            "offset": 0,
            "mode": "contain",
        },
        "{campaign_summary_picture}": {
            "replace_existing_picture": True,
            "offset": 0,
            "mode": "cover",
        },
    }

    def __init__(self, pptx_file: str) -> None:
        self.pptx_file = pptx_file
        self.prs = Presentation(pptx_file)
        self._slide_placeholders: Dict[int, set] = {}
        self._shape_channels: Dict[tuple, set] = {}
        self._channel_row_tops: Dict[int, Dict[str, list]] = {}

    def replace_placeholders(
        self,
        replacements: Dict[str, str],
        channels_present,
        output_pptx: str,
        image_placeholders: Optional[Dict[str, str]] = None,
    ) -> None:
        if isinstance(channels_present, bool):
            include_pin = channels_present
            channels_present = ["meta"]
            if include_pin:
                channels_present.append("pin")
        elif channels_present is None:
            channels_present = ["meta"]
        else:
            channels_present = list(channels_present)

        self._slide_placeholders = {}
        self._shape_channels = {}
        self._channel_row_tops = {}
        found_placeholders = set()
        display_lookup = {value: key for key, value in self._CHANNEL_DISPLAY.items()}
        image_replacements = {}
        if image_placeholders:
            image_replacements = {
                self._normalize_placeholder(key): value
                for key, value in image_placeholders.items()
                if value
            }

        for slide_index, slide in enumerate(self.prs.slides):
            slide_placeholders: set = set()
            for shape in slide.shapes:
                shape_channel_tags: set = set()
                if shape.has_text_frame:
                    image_replaced = False
                    for paragraph in shape.text_frame.paragraphs:
                        full_text = "".join(run.text for run in paragraph.runs)
                        image_hit = None
                        if image_replacements:
                            image_hit = self._find_image_placeholder(full_text, image_replacements)
                        if image_hit:
                            self._place_image(slide, shape, image_replacements[image_hit], image_hit)
                            found_placeholders.add(image_hit)
                            slide_placeholders.add(image_hit)
                            image_replaced = True
                            break
                        for token in self._PLACEHOLDER_PATTERN.findall(full_text):
                            normalized = self._normalize_placeholder(token)
                            channel = self._channel_from_placeholder(normalized)
                            if channel:
                                shape_channel_tags.add(channel)
                        new_text = self._apply_replacements(
                            full_text, replacements, found_placeholders, slide_placeholders
                        )
                        if new_text != full_text:
                            if not paragraph.runs:
                                paragraph.add_run()
                            paragraph.runs[0].text = new_text
                            for run in paragraph.runs[1:]:
                                run.text = ""
                        normalized_text = full_text.strip().lower()
                        if normalized_text in display_lookup:
                            shape_channel_tags.add(display_lookup[normalized_text])
                    if image_replaced:
                        continue

                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                cell_text = "".join(run.text for run in paragraph.runs)
                                for token in self._PLACEHOLDER_PATTERN.findall(cell_text):
                                    normalized = self._normalize_placeholder(token)
                                    channel = self._channel_from_placeholder(normalized)
                                    if channel:
                                        shape_channel_tags.add(channel)
                                new_text = self._apply_replacements(
                                    cell_text, replacements, found_placeholders, slide_placeholders
                                )
                                if new_text != cell_text:
                                    if not paragraph.runs:
                                        run = paragraph.add_run()
                                    else:
                                        run = paragraph.runs[0]
                                    run.text = new_text
                                    for extra_run in paragraph.runs[1:]:
                                        extra_run.text = ""
                                    for channel in self._detect_channel_from_text(new_text):
                                        shape_channel_tags.add(channel)
                                normalized_cell_text = cell_text.strip().lower()
                                if normalized_cell_text in display_lookup:
                                    shape_channel_tags.add(display_lookup[normalized_cell_text])
                if shape_channel_tags:
                    self._shape_channels[(slide_index, shape._element)] = shape_channel_tags
                    row_map = self._channel_row_tops.setdefault(slide_index, {})
                    for channel in shape_channel_tags:
                        row_map.setdefault(channel, []).append(shape.top)

            self._remove_placeholder_only_rows(slide)
            self._slide_placeholders[slide_index] = slide_placeholders

        self._prune_channel_slides(channels_present)
        self._refresh_summary_charts()

        self.prs.save(output_pptx)

        if not found_placeholders:
            print("⚠️ No placeholders were replaced. Verify placeholders in the PPT match the replacements.")

    # ------------------------------------------------------------------ #
    # Helpers
    # ------------------------------------------------------------------ #

    def _apply_replacements(
        self,
        text: str,
        replacements: Dict[str, str],
        global_found: set,
        slide_found: set,
    ) -> str:
        def replace_match(match: re.Match) -> str:
            raw = match.group(0)
            normalized = self._normalize_placeholder(raw)
            if normalized in replacements:
                global_found.add(normalized)
                slide_found.add(normalized)
                value = str(replacements[normalized])
                full_text = match.string
                prefix_char = None
                suffix_char = None

                idx = match.start() - 1
                while idx >= 0 and full_text[idx].isspace():
                    idx -= 1
                if idx >= 0:
                    prefix_char = full_text[idx]

                idx = match.end()
                length = len(full_text)
                while idx < length and full_text[idx].isspace():
                    idx += 1
                if idx < length:
                    suffix_char = full_text[idx]

                if prefix_char == "£":
                    value = value.lstrip(" £")
                if suffix_char == "%":
                    value = value.rstrip(" %")

                return value
            return raw

        return self._PLACEHOLDER_PATTERN.sub(replace_match, text)

    def _remove_placeholder_only_rows(self, slide) -> None:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                row_indexes = []
                for idx, row in enumerate(table.rows):
                    texts = [cell.text.strip() for cell in row.cells]
                    if texts and all(text.startswith("{") and text.endswith("}") for text in texts):
                        row_indexes.append(idx)
                for idx in reversed(row_indexes):
                    tbl = table._tbl
                    tr = tbl.tr_lst[idx]
                    tbl.remove(tr)

    def _prune_channel_slides(self, channels_present: Sequence[str]) -> None:
        channel_tokens = {
            "pin": ["{pin_", "{p_", "pinterest"],
            "tik": ["{tik_", "{t_", "tiktok"],
        }
        channel_tokens["meta"] = ["{meta_", "{m_", "meta"]

        present = set(channels_present)
        missing = set(channel_tokens.keys()) - present
        slide_indices_to_remove = []

        for index, slide in enumerate(self.prs.slides):
            # Remove rows or shapes within the slide that belong to absent channels
            if missing:
                self._remove_channel_specific_content(slide, index, missing, channel_tokens)

            placeholders = self._slide_placeholders.get(index, set())
            tags = {
                self._channel_from_placeholder(ph)
                for ph in placeholders
                if self._channel_from_placeholder(ph) is not None
            }

            if not tags:
                slide_text = self._collect_slide_text(slide).lower()
                for channel, tokens in channel_tokens.items():
                    if any(token in slide_text for token in tokens):
                        tags.add(channel)

            missing_tags = tags - present
            if tags and missing_tags == tags:
                slide_indices_to_remove.append(index)

        for index in reversed(slide_indices_to_remove):
            try:
                xml_slides = self.prs.slides._sldIdLst
                slide_id = xml_slides[index]
                xml_slides.remove(slide_id)
            except IndexError:
                continue

    def _remove_channel_specific_content(
        self,
        slide,
        slide_index: int,
        missing_channels: Iterable[str],
        channel_tokens: Dict[str, Iterable[str]],
    ) -> None:
        display_names = {ch: self._CHANNEL_DISPLAY.get(ch, ch).lower() for ch in channel_tokens}

        for shape in list(slide.shapes):
            shape_channels = self._shape_channels.get((slide_index, shape._element), set())
            if shape_channels and set(shape_channels).issubset(set(missing_channels)):
                self._remove_shape(slide, shape)
                continue

            if shape.has_table:
                table = shape.table
                rows_to_remove = []
                for idx, row in enumerate(table.rows):
                    row_texts = [cell.text for cell in row.cells]
                    if self._row_is_all_na(row_texts):
                        rows_to_remove.append(idx)
                        continue
                    first_cell_text = row_texts[0].strip().lower() if row_texts else ""
                    for channel in missing_channels:
                        display = display_names.get(channel)
                        if display and first_cell_text == display:
                            rows_to_remove.append(idx)
                            break
                for idx in reversed(sorted(set(rows_to_remove))):
                    tbl = table._tbl
                    tr = tbl.tr_lst[idx]
                    tbl.remove(tr)
                if len(table.rows) == 0 or all(self._row_is_all_na([cell.text for cell in row.cells]) for row in table.rows):
                    self._remove_shape(slide, shape)
            elif shape.has_text_frame:
                text = shape.text.strip().lower()
                if not text:
                    continue
                for channel in missing_channels:
                    display = display_names.get(channel)
                    if text == display or (
                        text in ("n/a", "na")
                    ):
                        self._remove_shape(slide, shape)
                        break
                    if display and display in text:
                        self._remove_shape(slide, shape)
                        break

        row_top_entries = self._channel_row_tops.get(slide_index, {})
        if not row_top_entries:
            return

        row_tops = {channel: min(tops) for channel, tops in row_top_entries.items() if tops}
        if not row_tops:
            return

        missing_in_rows = [ch for ch in missing_channels if ch in row_tops]
        if not missing_in_rows:
            return

        sorted_rows = sorted(row_tops.items(), key=lambda kv: kv[1])
        for channel, top in sorted(sorted(((ch, row_tops[ch]) for ch in missing_in_rows), key=lambda kv: kv[1])):
            next_tops = [t for c, t in sorted_rows if t > top and c not in missing_in_rows]
            if not next_tops:
                continue
            delta = min(next_tops) - top
            if delta <= 0:
                continue
            for shape in slide.shapes:
                if shape.top > top:
                    shape.top = max(shape.top - delta, 0)

    @staticmethod
    def _clean_na_token(text: str) -> str:
        cleaned = text.strip().lower()
        cleaned = cleaned.replace("£", "").replace("%", "")
        cleaned = cleaned.replace("\n", " ")
        cleaned = re.sub(r"\broas\b|\broi\b", "", cleaned)
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        return cleaned

    @staticmethod
    def _normalize_placeholder(raw: str) -> str:
        raw = raw.strip()
        if not raw.startswith("{") or not raw.endswith("}"):
            return raw
        body = raw[1:-1].replace(" ", "")
        return "{" + body + "}"

    @classmethod
    def _is_na_text(cls, text: str) -> bool:
        cleaned = cls._clean_na_token(text)
        return cleaned in {"", "n/a", "na"}

    @classmethod
    def _row_is_all_na(cls, row_texts: Iterable[str]) -> bool:
        texts = list(row_texts)
        if not texts:
            return True
        return all(cls._is_na_text(text) for text in texts)

    @staticmethod
    def _channel_from_placeholder(placeholder: str):
        norm = placeholder.strip("{}").lower()
        if norm.startswith(("meta_", "m_")):
            return "meta"
        if norm.startswith(("pin_", "p_")):
            return "pin"
        if norm.startswith(("tik_", "t_")):
            return "tik"
        return None

    def _detect_channel_from_text(self, text: str) -> Iterable[str]:
        channels = set()
        tokens = text.lower().replace("{", "").replace("}", "")
        if "meta" in tokens or "m_" in tokens:
            channels.add("meta")
        if "p_" in tokens or "pin_" in tokens or "pinterest" in tokens:
            channels.add("pin")
        if "tik" in tokens or "t_" in tokens or "tiktok" in tokens:
            channels.add("tik")
        return channels

    @staticmethod
    def _find_image_placeholder(text: str, image_replacements: Dict[str, str]) -> Optional[str]:
        if not text:
            return None
        normalized_text = "".join(text.split())
        for placeholder in image_replacements.keys():
            if placeholder in normalized_text or placeholder in text:
                return placeholder
        return None

    def _place_image(self, slide, shape, image_path: str, placeholder: str) -> None:
        if not image_path or not os.path.exists(image_path):
            return

        override = self._IMAGE_LAYOUT_OVERRIDES.get(placeholder, {})

        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE and not override:
            try:
                shape.insert_picture(image_path)
            except Exception:
                pass
            return

        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height

        def resolve_dimension(value, factor, shape_value, slide_value):
            if isinstance(value, str):
                if value == "slide":
                    return slide_value
            if value is not None:
                return value
            if isinstance(factor, (int, float)):
                return int(slide_value * factor)
            return shape_value

        target_picture = None
        if override.get("replace_existing_picture"):
            pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
            if pictures:
                target_picture = max(pictures, key=lambda sh: sh.width * sh.height)

        if target_picture is not None:
            left = override.get("left", target_picture.left)
            top = override.get("top", target_picture.top)
            width = resolve_dimension(
                override.get("width"),
                override.get("width_factor"),
                target_picture.width,
                slide_width,
            )
            height = resolve_dimension(
                override.get("height"),
                override.get("height_factor"),
                target_picture.height,
                slide_height,
            )
        else:
            left = override.get("left", shape.left)
            top = override.get("top", shape.top)
            width = resolve_dimension(
                override.get("width"),
                override.get("width_factor"),
                shape.width,
                slide_width,
            )
            height = resolve_dimension(
                override.get("height"),
                override.get("height_factor"),
                shape.height,
                slide_height,
            )

        if override.get("left") is None and target_picture is None and override.get("width_factor"):
            left = int((slide_width - width) / 2)
        if override.get("top") is None and target_picture is None and override.get("height_factor"):
            top = int((slide_height - height) / 2)

        if left is None:
            left = shape.left
        if top is None:
            top = shape.top
        if width is None:
            width = slide_width
        if height is None:
            height = slide_height

        offset_ratio = override.get("offset", 0.03 if width else 0)
        offset = int(width * offset_ratio) if isinstance(offset_ratio, (int, float)) else 0

        left = min(max(0, left), slide_width - width)
        top = min(max(0, top), slide_height - height)

        if target_picture is not None:
            self._remove_shape(slide, target_picture)
        self._remove_shape(slide, shape)

        try:
            picture = slide.shapes.add_picture(image_path, left + offset, top)
        except Exception:
            return

        if width and height and picture.width and picture.height:
            try:
                scale_w = width / picture.width
                scale_h = height / picture.height
                mode = override.get("mode", "cover")
                if mode == "contain":
                    scale = min(scale_w, scale_h)
                else:
                    scale = max(scale_w, scale_h)
            except ZeroDivisionError:
                scale = None
            if scale:
                picture.width = int(picture.width * scale)
                picture.height = int(picture.height * scale)
            picture.left = left + int((width - picture.width) / 2)
            picture.top = top + int((height - picture.height) / 2)

    @staticmethod
    def _remove_shape(slide, shape) -> None:
        slide.shapes._spTree.remove(shape._element)

    def _refresh_summary_charts(self) -> None:
        placeholder_key = "{estimated_versus_actual_performance_commentary}"
        normalized_key = self._normalize_placeholder(placeholder_key)

        target_index = None
        for idx, placeholders in self._slide_placeholders.items():
            if normalized_key in placeholders:
                target_index = idx
                break

        if target_index is None:
            if len(self.prs.slides) <= 8:
                return
            target_index = 8

        slide = self.prs.slides[target_index]

        for shape in list(slide.shapes):
            if shape.shape_type == 13:  # pictures
                slide.shapes._spTree.remove(shape._element)

        try:
            slide.shapes.add_picture("impressions_chart.png", Inches(0.7), Inches(2.9), height=Inches(3.2))
            slide.shapes.add_picture("clicks_chart.png", Inches(8.7), Inches(2.9), height=Inches(3.2))
        except FileNotFoundError:
            pass

    @staticmethod
    def _collect_slide_text(slide) -> str:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        return " ".join(texts)
