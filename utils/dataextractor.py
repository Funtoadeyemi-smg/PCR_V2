from __future__ import annotations

import json
import math
import os
import sys
from dataclasses import dataclass, asdict
from datetime import date, datetime
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

import numpy as np
import pandas as pd
import re
from utils._totals import coerce_total_row
from dotenv import load_dotenv
from openai import OpenAI
from pydantic import BaseModel, model_validator
from utils.charts import generate_summary_charts

load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client: Optional[OpenAI] = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None

sys.modules.setdefault("utils.dataextractor", sys.modules[__name__])

PRECISION = 2


# ---------------------------------------------------------------------------
# Pydantic validation
# ---------------------------------------------------------------------------

REQUIRED_COLUMNS = {
    "meta": {
        "Gross_Spend",
        "Impressions",
        "Clicks",
        "Reach",
        "Ad_Set_Reach",
        "Brand_Revenue",
        "Brand_Units",
        "Brand_Online_Revenue",
        "Brand_Instore_Revenue",
        "SKU_Revenue",
        "SKU_Units",
    },
    "pinterest": {
        "Spend in account currency",
        "Impressions",
        "Reach",
        "Pin clicks",
        "Total order value (Lead)",
        "Web order value (Lead)",
        "Offline order value (Lead)",
    },
    "tik_audience": {
        "Ad group name",
        "Cost",
        "Impressions",
        "Reach",
        "Frequency",
        "Clicks (destination)",
        "CTR (destination)",
        "Video views",
        "6-second video views",
    },
    "tik_ad": {
        "Ad name",
        "Cost",
        "Impressions",
        "Reach",
        "Frequency",
        "Clicks (destination)",
        "CTR (destination)",
        "Video views",
        "6-second video views",
    },
    "media_plan": {
        "Platform",
        "Estimated Impressions",
        "Estimated link clicks",
        "Estimated Frequency",
        "Estimated Reach",
        "Gross spend by channel / platform",
        "Estimated CTR",
        "Net CPM",
        "Flight duration",
    },
}


class DataFrameSchema(BaseModel):
    """Ensures required columns exist for a dataframe input."""

    name: str
    columns: List[str]

    @model_validator(mode="after")
    def validate_columns(self) -> "DataFrameSchema":
        required = REQUIRED_COLUMNS.get(self.name, set())
        available = {col.strip() for col in self.columns}
        missing = required - available
        if missing:
            raise ValueError(f"Missing required columns for {self.name}: {sorted(missing)}")
        return self


def validate_dataframe(name: str, df: pd.DataFrame) -> None:
    """Trigger a pydantic validation for dataframe columns."""
    DataFrameSchema(name=name, columns=list(map(str, df.columns)))


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def safe_div(numerator: float, denominator: float) -> float:
    if not denominator:
        return 0.0
    return numerator / denominator


def fmt_number(value: float, decimals: int = PRECISION) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "N/A"
    return f"{value:,.{decimals}f}"


def fmt_int(value: float) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "N/A"
    return f"{int(round(value)):,}"


def fmt_percent(value: float, decimals: int = 2) -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "N/A"
    return f"{value:.{decimals}f}%"


def fmt_currency(value: float, currency_symbol: str = "£") -> str:
    if value is None or (isinstance(value, float) and math.isnan(value)):
        return "N/A"
    return f"{currency_symbol}{value:,.2f}"


def coerce_float(value: object) -> float:
    if value is None:
        return 0.0
    if isinstance(value, (int, float, np.integer, np.floating)):
        if isinstance(value, (float, np.floating)) and math.isnan(float(value)):
            return 0.0
        return float(value)
    if isinstance(value, str):
        cleaned = value.replace(",", " ")
        match = re.search(r"[-+]?\d*\.?\d+", cleaned)
        if match:
            return float(match.group())
    return 0.0


def detect_tiktok_file_role(df: pd.DataFrame) -> str:
    if "Ad group name" in df.columns:
        return "audience"
    if "Ad name" in df.columns:
        return "ad"
    raise ValueError("Unable to determine TikTok file type. Expected 'Ad group name' or 'Ad name' column.")


def drop_total_rows(df: pd.DataFrame, key_column: str) -> pd.DataFrame:
    if key_column not in df.columns:
        return df
    mask = df[key_column].astype(str).str.lower().str.startswith("total")
    return df[~mask].copy()


# ---------------------------------------------------------------------------
# Dataclasses for structured handling
# ---------------------------------------------------------------------------


@dataclass
class ChannelSummary:
    channel: str
    gross_spend: float
    net_spend: float
    impressions: float
    reach: float
    clicks: float
    ctr: float
    net_cpm: float
    frequency: float
    brand_revenue: float
    brand_roas: float
    brand_roi: float
    brand_units: float
    brand_online_revenue: float
    brand_instore_revenue: float
    brand_online_units: float
    brand_instore_units: float
    fsku_revenue: float
    fsku_roas: float
    fsku_roi: float
    fsku_units: float
    fsku_online_revenue: float
    fsku_instore_revenue: float


@dataclass
class AudiencePerformance:
    name: str
    net_spend: float
    impressions: float
    reach: float
    frequency: float
    clicks: float
    ctr: float
    net_cpm: float
    revenue: float
    roas: float
    roi: float


@dataclass
class CreativePerformance:
    name: str
    net_spend: float
    impressions: float
    reach: float
    frequency: float
    clicks: float
    ctr: float
    net_cpm: float
    revenue: float
    roas: float


# ---------------------------------------------------------------------------
# Main extractor
# ---------------------------------------------------------------------------


class ConsolidatedDataExtractor:
    COMMENTARY_KEYS_IN_ORDER = [
        "overall_campaign_performance_commentary",
        "overall_engagement_performance_commentary",
        "overall_sales_revenue_ROI_performance_commentary",
        "overall_reach_performance_commentary",
        "estimated_versus_actual_performance_commentary",
        "channel_performance_commentary",
        "brand_and_fsku_roas_performance_wrt_overall_traffic_performance_commentary",
        "percentage_of_instore_online_sales_commentary",
        "meta_two_audiences_with_strongest_roas_commentary",
        "meta_intuitive_commentary_about_audience_performance_commentary",
        "meta_audience_performance_commentary",
        "meta_ad_with_strongest_roas_and_impressions_commentary",
        "meta_ad_with_weakest_roas_and_impressions_commentary",
        "tik_two_audiences_with_strongest_roas_commentary",
        "tik_intuitive_commentary_about_audience_performance_commentary",
        "tik_audience_performance_commentary",
        "tik_ad_with_strongest_roas_and_impressions_commentary",
        "tik_ad_with_weakest_roas_and_impressions_commentary",
        "pin_two_audiences_with_strongest_roas_commentary",
        "pin_intuitive_commentary_about_audience_performance_commentary",
        "pin_audience_performance_commentary",
        "pin_ad_with_strongest_roas_and_impressions_commentary",
        "pin_ad_with_weakest_roas_and_impressions_commentary",
        "how_objective_2_was_met_commentary",
        "how_objective_1_was_met_commentary",
    ]

    CHANNEL_COMMENTARY_KEYS = {
        "meta": [
            "meta_two_audiences_with_strongest_roas_commentary",
            "meta_intuitive_commentary_about_audience_performance_commentary",
            "meta_audience_performance_commentary",
            "meta_ad_with_strongest_roas_and_impressions_commentary",
            "meta_ad_with_weakest_roas_and_impressions_commentary",
        ],
        "tik": [
            "tik_two_audiences_with_strongest_roas_commentary",
            "tik_intuitive_commentary_about_audience_performance_commentary",
            "tik_audience_performance_commentary",
            "tik_ad_with_strongest_roas_and_impressions_commentary",
            "tik_ad_with_weakest_roas_and_impressions_commentary",
        ],
        "pin": [
            "pin_two_audiences_with_strongest_roas_commentary",
            "pin_intuitive_commentary_about_audience_performance_commentary",
            "pin_audience_performance_commentary",
            "pin_ad_with_strongest_roas_and_impressions_commentary",
            "pin_ad_with_weakest_roas_and_impressions_commentary",
        ],
    }

    def __init__(
        self,
        meta_file: str,
        pinterest_file: Optional[str] = None,
        tiktok_files: Optional[Sequence[str]] = None,
        media_plan_file: Optional[str] = None,
        prompt_template: Optional[str] = None,
        template_path: Optional[str] = None,
        manual_campaign_objective: Optional[str] = None,
        manual_primary_kpis: Optional[str] = None,
        manual_secondary_kpis: Optional[str] = None,
    ) -> None:
        self.meta_file = meta_file
        self.pinterest_file = pinterest_file
        self.tiktok_files = list(tiktok_files or [])
        self.media_plan_file = media_plan_file
        self.prompt_template = prompt_template
        self.template_path = template_path or "consolidated_template.pptx"
        self.manual_campaign_objective = (manual_campaign_objective or "").strip()
        self.manual_primary_kpis = (manual_primary_kpis or "").strip()
        self.manual_secondary_kpis = (manual_secondary_kpis or "").strip()

        self.channels_present: List[str] = ["meta"]

    # ------------------------------------------------------------------ #
    # Public API
    # ------------------------------------------------------------------ #

    def extract_values(self) -> Tuple[Dict[str, str], List[str]]:
        meta_df, meta_conversion_frames = self._load_meta_frames(self.meta_file)
        validate_dataframe("meta", meta_df)

        pinterest_df = None
        if self.pinterest_file:
            pinterest_df = self._load_generic_dataframe(self.pinterest_file)
            validate_dataframe("pinterest", pinterest_df)
            self.channels_present.append("pin")

        tiktok_audience_df, tiktok_ad_df = self._load_tiktok_frames(self.tiktok_files)
        if tiktok_audience_df is not None or tiktok_ad_df is not None:
            self.channels_present.append("tik")

        media_plan_df = None
        if self.media_plan_file:
            media_plan_df = self._load_media_plan(self.media_plan_file)
            validate_dataframe("media_plan", media_plan_df)

        meta_summary, meta_audiences, meta_ads = self._build_meta_summary(meta_df, meta_conversion_frames)

        pinterest_summary = pinterest_audiences = pinterest_ads = None
        if pinterest_df is not None:
            pinterest_summary, pinterest_audiences, pinterest_ads = self._build_pinterest_summary(pinterest_df)

        tiktok_summary = tiktok_audiences = tiktok_ads = None
        if tiktok_audience_df is not None or tiktok_ad_df is not None:
            tiktok_summary, tiktok_audiences, tiktok_ads = self._build_tiktok_summary(
                tiktok_audience_df, tiktok_ad_df
            )

        estimates = self._build_estimates(media_plan_df)
        overall = self._build_overall_summary(meta_summary, pinterest_summary, tiktok_summary)
        est_totals = self._aggregate_estimates(estimates)

        channel_context = {
            "meta": {
                "summary": meta_summary,
                "audiences": meta_audiences or [],
                "creatives": meta_ads or [],
            }
        }
        if pinterest_summary:
            channel_context["pin"] = {
                "summary": pinterest_summary,
                "audiences": pinterest_audiences or [],
                "creatives": pinterest_ads or [],
            }
        if tiktok_summary:
            channel_context["tik"] = {
                "summary": tiktok_summary,
                "audiences": tiktok_audiences or [],
                "creatives": tiktok_ads or [],
            }

        replacements: Dict[str, str] = {}
        objectives = self._extract_objectives(meta_df)

        replacements.update(self._basic_identity_fields(meta_df, media_plan_df, objectives))
        replacements.update(self._overall_metrics(overall, estimates, est_totals))

        replacements.update(self._channel_placeholders("meta", meta_summary, estimates.get("meta"), meta_audiences, meta_ads))

        if pinterest_summary:
            replacements.update(
                self._channel_placeholders("pin", pinterest_summary, estimates.get("pin"), pinterest_audiences, pinterest_ads)
            )
        else:
            replacements.update(self._empty_channel_placeholders("pin"))

        if tiktok_summary:
            replacements.update(
                self._channel_placeholders("tik", tiktok_summary, estimates.get("tik"), tiktok_audiences, tiktok_ads)
            )
        else:
            replacements.update(self._empty_channel_placeholders("tik"))

        replacements.update(self._audience_tables(meta_audiences, pinterest_audiences, tiktok_audiences))

        replacements.update(self._creative_placeholders("meta", meta_ads))
        replacements.update(self._creative_placeholders("pin", pinterest_ads))
        replacements.update(self._creative_placeholders("tik", tiktok_ads))

        online_replacements, online_metrics = self._online_instore_metrics(meta_summary, pinterest_summary)
        replacements.update(online_replacements)

        replacements.update(
            self._generate_commentaries(
                overall,
                estimates,
                channel_context,
                online_metrics,
                objectives,
            )
        )

        replacements = self._fill_missing_placeholders(replacements)

        self._generate_summary_charts(overall, est_totals)

        return replacements, self.channels_present

    # ------------------------------------------------------------------ #
    # Loading helpers
    # ------------------------------------------------------------------ #

    def _load_generic_dataframe(self, path: str) -> pd.DataFrame:
        suffix = Path(path).suffix.lower()
        if suffix in (".csv", ".txt"):
            return pd.read_csv(path)
        return pd.read_excel(path)

    def _load_media_plan(self, path: str) -> pd.DataFrame:
        required_columns = {col.lower() for col in REQUIRED_COLUMNS["media_plan"]}

        try:
            preview = pd.read_excel(path, header=None, nrows=20)
        except Exception:
            # Fall back to the default reader; validation will surface any issues.
            df = pd.read_excel(path)
        else:
            header_row = 0
            for idx, row in preview.iterrows():
                normalized = {
                    str(value).strip().lower()
                    for value in row.tolist()
                    if pd.notna(value) and str(value).strip()
                }
                if required_columns.issubset(normalized):
                    header_row = idx
                    break
            df = pd.read_excel(path, header=header_row)

        df.columns = [str(col).strip() for col in df.columns]
        df = df.dropna(how="all")
        return df

    def _load_meta_frames(self, path: str) -> Tuple[pd.DataFrame, List[pd.DataFrame]]:
        suffix = Path(path).suffix.lower()
        if suffix in (".csv", ".txt"):
            return pd.read_csv(path), []

        workbook = pd.ExcelFile(path)
        main_sheet = workbook.sheet_names[0]
        main_df = workbook.parse(main_sheet)

        conversion_frames: List[pd.DataFrame] = []
        for sheet in workbook.sheet_names:
            normalized = sheet.strip().lower().replace(" ", "")
            if "conv" in normalized:
                try:
                    frame = workbook.parse(sheet)
                    conversion_frames.append(frame)
                except Exception:
                    continue

        return main_df, conversion_frames

    def _load_tiktok_frames(
        self, files: Sequence[str]
    ) -> Tuple[Optional[pd.DataFrame], Optional[pd.DataFrame]]:
        audience_df = None
        ad_df = None

        for file_path in files:
            df = self._load_generic_dataframe(file_path)
            role = detect_tiktok_file_role(df)
            if role == "audience":
                validate_dataframe("tik_audience", df)
                audience_df = df
            else:
                validate_dataframe("tik_ad", df)
                ad_df = df

        return audience_df, ad_df

    # ------------------------------------------------------------------ #
    # Metric builders
    # ------------------------------------------------------------------ #

    def _extract_meta_conversion_metrics(self, frames: Sequence[pd.DataFrame]) -> Dict[str, float]:
        if not frames:
            return {}

        processed: List[pd.DataFrame] = []
        for frame in frames:
            if frame is None or frame.empty:
                continue
            df = frame.copy()
            df.columns = [str(col).strip() for col in df.columns]
            required = {"Level", "Hier", "Type", "Units", "Sales"}
            if not required.issubset(df.columns):
                continue
            processed.append(df)

        if not processed:
            return {}

        combined = pd.concat(processed, ignore_index=True)
        combined["Level"] = combined["Level"].astype(str)
        combined["Hier"] = combined["Hier"].astype(str)
        combined["Type"] = combined["Type"].astype(str)
        combined["Units"] = pd.to_numeric(combined["Units"], errors="coerce").fillna(0.0)
        combined["Sales"] = pd.to_numeric(combined["Sales"], errors="coerce").fillna(0.0)

        brand_mask = (combined["Level"].str.upper() == "BRAND") & (combined["Hier"].str.lower() == "campaign")
        brand_conv = combined[brand_mask]

        sku_mask = (combined["Level"].str.upper() == "SKU") & (combined["Hier"].str.lower() == "campaign")
        sku_conv = combined[sku_mask]

        metrics: Dict[str, float] = {}
        if not brand_conv.empty:
            brand_grouped = brand_conv.groupby(brand_conv["Type"].str.lower())
            unit_totals = brand_grouped["Units"].sum()
            sales_totals = brand_grouped["Sales"].sum()
            metrics["brand_online_units"] = float(unit_totals.get("online", 0.0))
            metrics["brand_instore_units"] = float(unit_totals.get("offline", 0.0))
            metrics["brand_online_revenue"] = float(sales_totals.get("online", 0.0))
            metrics["brand_instore_revenue"] = float(sales_totals.get("offline", 0.0))

        return metrics

    def _build_meta_summary(
        self,
        df: pd.DataFrame,
        conversion_frames: Optional[Sequence[pd.DataFrame]] = None,
    ) -> Tuple[ChannelSummary, List[AudiencePerformance], List[CreativePerformance]]:
        conversion_metrics = self._extract_meta_conversion_metrics(conversion_frames or [])
        df = df.copy()
        df_with_totals = coerce_total_row(df, key_column="Ad_Set_Name", numeric_columns=[
            "Gross_Spend",
            "Net_Spend",
            "Impressions",
            "Reach",
            "Ad_Set_Reach",
            "Clicks",
            "CTR",
            "CPM",
            "Brand_Revenue",
            "Brand_Units",
            "Brand_Online_Revenue",
            "Brand_Instore_Revenue",
            "SKU_Revenue",
            "SKU_Units",
            "SKU_Online_Revenue",
            "SKU_Instore_Revenue",
        ])

        mask_total = (
            df_with_totals["Ad_Set_Name"].astype(str).str.startswith("Total", na=False)
            if "Ad_Set_Name" in df_with_totals.columns
            else pd.Series(False, index=df_with_totals.index)
        )
        totals_row = df_with_totals[mask_total].iloc[-1] if mask_total.any() else None
        df = df_with_totals[~mask_total].copy() if mask_total.any() else df_with_totals.copy()

        df["Gross_Spend"] = df["Gross_Spend"].fillna(0)
        df["Net_Spend"] = df.get("Net_Spend", df["Gross_Spend"] * 0.7).fillna(df["Gross_Spend"] * 0.7)
        df["Impressions"] = df["Impressions"].fillna(0)
        df["Reach"] = df["Reach"].fillna(df["Ad_Set_Reach"].fillna(0))
        df["Clicks"] = df["Clicks"].fillna(0)

        if totals_row is not None:
            gross_spend = coerce_float(totals_row.get("Gross_Spend"))
            net_spend = coerce_float(totals_row.get("Net_Spend"))
            impressions = coerce_float(totals_row.get("Impressions"))
            reach = coerce_float(totals_row.get("Reach")) or coerce_float(totals_row.get("Ad_Set_Reach"))
            clicks = coerce_float(totals_row.get("Clicks"))
            brand_revenue = coerce_float(totals_row.get("Brand_Revenue"))
            brand_units = coerce_float(totals_row.get("Brand_Units"))
            brand_online = coerce_float(totals_row.get("Brand_Online_Revenue"))
            brand_instore = coerce_float(totals_row.get("Brand_Instore_Revenue"))
            fsku_revenue = coerce_float(totals_row.get("SKU_Revenue"))
            fsku_units = coerce_float(totals_row.get("SKU_Units"))
            fsku_online = coerce_float(totals_row.get("SKU_Online_Revenue"))
            fsku_instore = coerce_float(totals_row.get("SKU_Instore_Revenue"))
        else:
            gross_spend = float(df["Gross_Spend"].sum())
            net_spend = float(df["Net_Spend"].sum())
            impressions = float(df["Impressions"].sum())
            reach = float(df["Reach"].sum())
            clicks = float(df["Clicks"].sum())
            brand_revenue = float(df["Brand_Revenue"].fillna(0).sum())
            brand_units = float(df["Brand_Units"].fillna(0).sum())
            brand_online = float(df["Brand_Online_Revenue"].fillna(0).sum())
            brand_instore = float(df["Brand_Instore_Revenue"].fillna(0).sum())
            fsku_revenue = float(df["SKU_Revenue"].fillna(0).sum())
            fsku_units = float(df["SKU_Units"].fillna(0).sum())
            fsku_online = float(df.get("SKU_Online_Revenue", pd.Series([0])).fillna(0).sum())
            fsku_instore = float(df.get("SKU_Instore_Revenue", pd.Series([0])).fillna(0).sum())

        brand_online_units = 0.0
        brand_instore_units = 0.0
        has_brand_unit_split = False

        if conversion_metrics:
            conv_online_units = conversion_metrics.get("brand_online_units")
            conv_instore_units = conversion_metrics.get("brand_instore_units")
            if conv_online_units is not None:
                brand_online_units = conv_online_units
                has_brand_unit_split = True
            if conv_instore_units is not None:
                brand_instore_units = conv_instore_units
                has_brand_unit_split = True
            if conv_online_units is not None and conv_instore_units is not None:
                brand_units = brand_online_units + brand_instore_units

            conv_online_revenue = conversion_metrics.get("brand_online_revenue")
            conv_instore_revenue = conversion_metrics.get("brand_instore_revenue")
            if conv_online_revenue is not None:
                brand_online = conv_online_revenue
            if conv_instore_revenue is not None:
                brand_instore = conv_instore_revenue
            if conv_online_revenue is not None or conv_instore_revenue is not None:
                brand_revenue = brand_online + brand_instore

        if not has_brand_unit_split:
            if brand_revenue:
                online_share = safe_div(brand_online, brand_revenue)
                brand_online_units = brand_units * online_share
                brand_instore_units = max(0.0, brand_units - brand_online_units)
            else:
                brand_online_units = brand_units
                brand_instore_units = 0.0

        if not net_spend and gross_spend:
            net_spend = gross_spend * 0.7

        ctr = safe_div(clicks, impressions) * 100
        net_cpm = safe_div(net_spend, impressions) * 1000
        frequency = safe_div(impressions, reach)

        brand_roas = safe_div(brand_revenue, net_spend)
        brand_roi = safe_div(brand_revenue, gross_spend)

        fsku_roas = safe_div(fsku_revenue, net_spend)
        fsku_roi = safe_div(fsku_revenue, gross_spend)

        summary = ChannelSummary(
            channel="meta",
            gross_spend=gross_spend,
            net_spend=net_spend,
            impressions=impressions,
            reach=reach,
            clicks=clicks,
            ctr=ctr,
            net_cpm=net_cpm,
            frequency=frequency,
            brand_revenue=brand_revenue,
            brand_roas=brand_roas,
            brand_roi=brand_roi,
            brand_units=brand_units,
            brand_online_revenue=brand_online,
            brand_instore_revenue=brand_instore,
            brand_online_units=brand_online_units,
            brand_instore_units=brand_instore_units,
            fsku_revenue=fsku_revenue,
            fsku_roas=fsku_roas,
            fsku_roi=fsku_roi,
            fsku_units=fsku_units,
            fsku_online_revenue=fsku_online,
            fsku_instore_revenue=fsku_instore,
        )

        audience_df = (
            drop_total_rows(df, "Ad_Set_Name")
            .groupby("Ad_Set_Name", dropna=False)
            .agg(
                Net_Spend=("Net_Spend", "sum"),
                Impressions=("Impressions", "sum"),
                Reach=("Reach", "sum"),
                Clicks=("Clicks", "sum"),
                Brand_Revenue=("Brand_Revenue", "sum"),
            )
            .reset_index()
        )
        audience_df = audience_df[
            ~audience_df["Ad_Set_Name"].astype(str).str.lower().str.startswith("total")
        ]
        audience_df["CTR"] = audience_df["Clicks"].divide(audience_df["Impressions"].replace({0: np.nan})) * 100
        audience_df["Net_CPM"] = audience_df["Net_Spend"].divide(audience_df["Impressions"].replace({0: np.nan})) * 1000
        audience_df["Frequency"] = audience_df["Impressions"].divide(audience_df["Reach"].replace({0: np.nan}))
        audience_df["ROAS"] = audience_df["Brand_Revenue"].divide(audience_df["Net_Spend"].replace({0: np.nan}))
        audience_df["ROI"] = audience_df["Brand_Revenue"].divide(
            audience_df["Net_Spend"].replace({0: np.nan}) / 0.7
        )

        audiences = [
            AudiencePerformance(
                name=row["Ad_Set_Name"],
                net_spend=float(row["Net_Spend"]),
                impressions=float(row["Impressions"]),
                reach=float(row["Reach"]),
                frequency=float(row["Frequency"]) if not np.isnan(row["Frequency"]) else 0.0,
                clicks=float(row["Clicks"]),
                ctr=float(row["CTR"]) if not np.isnan(row["CTR"]) else 0.0,
                net_cpm=float(row["Net_CPM"]) if not np.isnan(row["Net_CPM"]) else 0.0,
                revenue=float(row["Brand_Revenue"]),
                roas=float(row["ROAS"]) if not np.isnan(row["ROAS"]) else 0.0,
                roi=float(row["ROI"]) if not np.isnan(row["ROI"]) else 0.0,
            )
            for _, row in audience_df.sort_values("Net_Spend", ascending=False).head(12).iterrows()
        ]

        ad_df = (
            drop_total_rows(df, "Ad")
            .groupby("Ad", dropna=False)
            .agg(
                Net_Spend=("Net_Spend", "sum"),
                Impressions=("Impressions", "sum"),
                Reach=("Reach", "sum"),
                Clicks=("Clicks", "sum"),
                Brand_Revenue=("Brand_Revenue", "sum"),
            )
            .reset_index()
        )
        ad_df = ad_df[~ad_df["Ad"].astype(str).str.lower().str.startswith("total")]
        ad_df["CTR"] = ad_df["Clicks"].divide(ad_df["Impressions"].replace({0: np.nan})) * 100
        ad_df["Net_CPM"] = ad_df["Net_Spend"].divide(ad_df["Impressions"].replace({0: np.nan})) * 1000
        ad_df["Frequency"] = ad_df["Impressions"].divide(ad_df["Reach"].replace({0: np.nan}))
        ad_df["ROAS"] = ad_df["Brand_Revenue"].divide(ad_df["Net_Spend"].replace({0: np.nan}))

        ads = [
            CreativePerformance(
                name=row["Ad"],
                net_spend=float(row["Net_Spend"]),
                impressions=float(row["Impressions"]),
                reach=float(row["Reach"]),
                frequency=float(row["Frequency"]) if not np.isnan(row["Frequency"]) else 0.0,
                clicks=float(row["Clicks"]),
                ctr=float(row["CTR"]) if not np.isnan(row["CTR"]) else 0.0,
                net_cpm=float(row["Net_CPM"]) if not np.isnan(row["Net_CPM"]) else 0.0,
                revenue=float(row["Brand_Revenue"]),
                roas=float(row["ROAS"]) if not np.isnan(row["ROAS"]) else 0.0,
            )
            for _, row in ad_df.iterrows()
        ]

        return summary, audiences, ads

    def _build_pinterest_summary(
        self, df: pd.DataFrame
    ) -> Tuple[ChannelSummary, List[AudiencePerformance], List[CreativePerformance]]:
        df = df.copy()
        df["Spend in account currency"] = df["Spend in account currency"].fillna(0)
        df["Impressions"] = df["Impressions"].fillna(0)
        df["Reach"] = df["Reach"].fillna(0)
        df["Pin clicks"] = df["Pin clicks"].fillna(0)
        df["Total order value (Lead)"] = df["Total order value (Lead)"].fillna(0)

        net_spend = float(df["Spend in account currency"].sum())
        gross_spend = net_spend / 0.7
        impressions = float(df["Impressions"].sum())
        reach = float(df["Reach"].sum())
        clicks = float(df["Pin clicks"].sum())
        ctr = safe_div(clicks, impressions) * 100
        net_cpm = safe_div(net_spend, impressions) * 1000
        frequency = safe_div(impressions, reach)

        brand_revenue = float(df["Total order value (Lead)"].sum())
        brand_roas = safe_div(brand_revenue, net_spend)
        brand_roi = safe_div(brand_revenue, gross_spend)

        web_revenue = float(df["Web order value (Lead)"].fillna(0).sum())
        offline_revenue = float(df["Offline order value (Lead)"].fillna(0).sum())

        summary = ChannelSummary(
            channel="pin",
            gross_spend=gross_spend,
            net_spend=net_spend,
            impressions=impressions,
            reach=reach,
            clicks=clicks,
            ctr=ctr,
            net_cpm=net_cpm,
            frequency=frequency,
            brand_revenue=brand_revenue,
            brand_roas=brand_roas,
            brand_roi=brand_roi,
            brand_units=float(df.get("Total conversions (Lead)", pd.Series([0])).fillna(0).sum()),
            brand_online_revenue=web_revenue,
            brand_instore_revenue=offline_revenue,
            brand_online_units=0.0,
            brand_instore_units=0.0,
            fsku_revenue=brand_revenue,
            fsku_roas=brand_roas,
            fsku_roi=brand_roi,
            fsku_units=float(df.get("Total conversions (Lead)", pd.Series([0])).fillna(0).sum()),
            fsku_online_revenue=web_revenue,
            fsku_instore_revenue=offline_revenue,
        )

        audience_df = (
            drop_total_rows(df, "Ad group name")
            .groupby("Ad group name", dropna=False)
            .agg(
                Net_Spend=("Spend in account currency", "sum"),
                Impressions=("Impressions", "sum"),
                Reach=("Reach", "sum"),
                Clicks=("Pin clicks", "sum"),
                Brand_Revenue=("Total order value (Lead)", "sum"),
            )
            .reset_index()
        )
        audience_df["CTR"] = audience_df["Clicks"].divide(audience_df["Impressions"].replace({0: np.nan})) * 100
        audience_df["Net_CPM"] = audience_df["Net_Spend"].divide(audience_df["Impressions"].replace({0: np.nan})) * 1000
        audience_df["Frequency"] = audience_df["Impressions"].divide(audience_df["Reach"].replace({0: np.nan}))
        audience_df["ROAS"] = audience_df["Brand_Revenue"].divide(audience_df["Net_Spend"].replace({0: np.nan}))
        audience_df["ROI"] = audience_df["Brand_Revenue"].divide(
            (audience_df["Net_Spend"] / 0.7).replace({0: np.nan})
        )

        audiences = [
            AudiencePerformance(
                name=row["Ad group name"],
                net_spend=float(row["Net_Spend"]),
                impressions=float(row["Impressions"]),
                reach=float(row["Reach"]),
                frequency=float(row["Frequency"]) if not np.isnan(row["Frequency"]) else 0.0,
                clicks=float(row["Clicks"]),
                ctr=float(row["CTR"]) if not np.isnan(row["CTR"]) else 0.0,
                net_cpm=float(row["Net_CPM"]) if not np.isnan(row["Net_CPM"]) else 0.0,
                revenue=float(row["Brand_Revenue"]),
                roas=float(row["ROAS"]) if not np.isnan(row["ROAS"]) else 0.0,
                roi=float(row["ROI"]) if not np.isnan(row["ROI"]) else 0.0,
            )
            for _, row in audience_df.sort_values("Net_Spend", ascending=False).head(12).iterrows()
        ]

        ad_df = (
            drop_total_rows(df, "Ad name")
            .groupby("Ad name", dropna=False)
            .agg(
                Net_Spend=("Spend in account currency", "sum"),
                Impressions=("Impressions", "sum"),
                Reach=("Reach", "sum"),
                Clicks=("Pin clicks", "sum"),
                Brand_Revenue=("Total order value (Lead)", "sum"),
            )
            .reset_index()
        )
        ad_df["CTR"] = ad_df["Clicks"].divide(ad_df["Impressions"].replace({0: np.nan})) * 100
        ad_df["Net_CPM"] = ad_df["Net_Spend"].divide(ad_df["Impressions"].replace({0: np.nan})) * 1000
        ad_df["Frequency"] = ad_df["Impressions"].divide(ad_df["Reach"].replace({0: np.nan}))
        ad_df["ROAS"] = ad_df["Brand_Revenue"].divide(ad_df["Net_Spend"].replace({0: np.nan}))

        ads = [
            CreativePerformance(
                name=row["Ad name"],
                net_spend=float(row["Net_Spend"]),
                impressions=float(row["Impressions"]),
                reach=float(row["Reach"]),
                frequency=float(row["Frequency"]) if not np.isnan(row["Frequency"]) else 0.0,
                clicks=float(row["Clicks"]),
                ctr=float(row["CTR"]) if not np.isnan(row["CTR"]) else 0.0,
                net_cpm=float(row["Net_CPM"]) if not np.isnan(row["Net_CPM"]) else 0.0,
                revenue=float(row["Brand_Revenue"]),
                roas=float(row["ROAS"]) if not np.isnan(row["ROAS"]) else 0.0,
            )
            for _, row in ad_df.iterrows()
        ]

        return summary, audiences, ads

    def _build_tiktok_summary(
        self, audience_df: Optional[pd.DataFrame], ad_df: Optional[pd.DataFrame]
    ) -> Tuple[ChannelSummary, List[AudiencePerformance], List[CreativePerformance]]:
        if audience_df is None and ad_df is None:
            raise ValueError("TikTok data requires at least an audience or ad level file.")

        if audience_df is not None:
            totals = coerce_total_row(
                audience_df,
                key_column="Ad group name",
                numeric_columns=[
                    "Cost",
                    "Impressions",
                    "Clicks (destination)",
                    "Video views",
                    "6-second video views",
                ],
            ).iloc[0]
            net_spend = float(totals.get("Cost", 0))
            impressions = float(totals.get("Impressions", 0))
            reach = float(totals.get("Reach", 0))
            clicks = float(totals.get("Clicks (destination)", 0))
            frequency = float(totals.get("Frequency", 0))
        else:
            base_df = ad_df.copy()
            base_df["Cost"] = base_df["Cost"].fillna(0)
            base_df["Impressions"] = base_df["Impressions"].fillna(0)
            base_df["Reach"] = base_df["Reach"].fillna(0)
            base_df["Frequency"] = base_df["Frequency"].fillna(0)
            base_df["Clicks (destination)"] = base_df["Clicks (destination)"].fillna(0)
            net_spend = float(base_df["Cost"].sum())
            impressions = float(base_df["Impressions"].sum())
            reach = float(base_df["Reach"].sum())
            clicks = float(base_df["Clicks (destination)"].sum())
            frequency = safe_div(impressions, reach)

        gross_spend = net_spend / 0.7 if net_spend else 0.0
        ctr = safe_div(clicks, impressions) * 100 if impressions else 0.0
        net_cpm = safe_div(net_spend, impressions) * 1000 if impressions else 0.0

        summary = ChannelSummary(
            channel="tik",
            gross_spend=gross_spend,
            net_spend=net_spend,
            impressions=impressions,
            reach=reach,
            clicks=clicks,
            ctr=ctr,
            net_cpm=net_cpm,
            frequency=frequency,
            brand_revenue=0.0,
            brand_roas=0.0,
            brand_roi=0.0,
            brand_units=0.0,
            brand_online_revenue=0.0,
            brand_instore_revenue=0.0,
            brand_online_units=0.0,
            brand_instore_units=0.0,
            fsku_revenue=0.0,
            fsku_roas=0.0,
            fsku_roi=0.0,
            fsku_units=0.0,
            fsku_online_revenue=0.0,
            fsku_instore_revenue=0.0,
        )

        audiences: List[AudiencePerformance] = []
        if audience_df is not None:
            filtered_audience = drop_total_rows(audience_df.copy(), "Ad group name")
            filtered_audience["CTR"] = filtered_audience["Clicks (destination)"].divide(
                filtered_audience["Impressions"].replace({0: np.nan})
            ) * 100
            filtered_audience["Net_CPM"] = filtered_audience["Cost"].divide(
                filtered_audience["Impressions"].replace({0: np.nan})
            ) * 1000
            filtered_audience["ROAS"] = 0.0
            filtered_audience["ROI"] = 0.0

            audiences = [
                AudiencePerformance(
                    name=row["Ad group name"],
                    net_spend=float(row["Cost"]),
                    impressions=float(row["Impressions"]),
                    reach=float(row["Reach"]),
                    frequency=float(row["Frequency"]) if not np.isnan(row["Frequency"]) else 0.0,
                    clicks=float(row["Clicks (destination)"]),
                    ctr=float(row["CTR"]) if not np.isnan(row["CTR"]) else 0.0,
                    net_cpm=float(row["Net_CPM"]) if not np.isnan(row["Net_CPM"]) else 0.0,
                    revenue=0.0,
                    roas=0.0,
                    roi=0.0,
                )
                for _, row in filtered_audience.sort_values("Cost", ascending=False).head(12).iterrows()
            ]

        creatives: List[CreativePerformance] = []
        if ad_df is not None:
            filtered_ad = drop_total_rows(ad_df.copy(), "Ad name")
            filtered_ad["CTR"] = filtered_ad["Clicks (destination)"].divide(
                filtered_ad["Impressions"].replace({0: np.nan})
            ) * 100
            filtered_ad["Net_CPM"] = filtered_ad["Cost"].divide(
                filtered_ad["Impressions"].replace({0: np.nan})
            ) * 1000

            creatives = [
                CreativePerformance(
                    name=row["Ad name"],
                    net_spend=float(row["Cost"]),
                    impressions=float(row["Impressions"]),
                    reach=float(row["Reach"]),
                    frequency=float(row["Frequency"]) if not np.isnan(row["Frequency"]) else 0.0,
                    clicks=float(row["Clicks (destination)"]),
                    ctr=float(row["CTR"]) if not np.isnan(row["CTR"]) else 0.0,
                    net_cpm=float(row["Net_CPM"]) if not np.isnan(row["Net_CPM"]) else 0.0,
                    revenue=0.0,
                    roas=0.0,
                )
                for _, row in filtered_ad.iterrows()
            ]

        return summary, audiences, creatives

    def _build_estimates(self, df: Optional[pd.DataFrame]) -> Dict[str, Dict[str, float]]:
        if df is None:
            return {}

        df = df.copy()
        original_platform = df["Platform"].copy()
        df["Platform"] = df["Platform"].fillna(method="ffill")

        selector_cols = [
            "Format & placement",
            "Targeting segment",
            "Audience Data Source",
            "Channel",
        ]
        channel_total_mask = original_platform.isna()
        for col in selector_cols:
            if col in df.columns:
                channel_total_mask &= df[col].isna()
        channel_total_mask &= df["Gross spend by channel / platform"].notna()
        df["_is_channel_total"] = channel_total_mask

        grouped = df.groupby("Platform")

        estimates: Dict[str, Dict[str, float]] = {}
        for platform, group in grouped:
            key = None
            if "Meta" in platform:
                key = "meta"
            elif "Pinterest" in platform:
                key = "pin"
            elif "TikTok" in platform:
                key = "tik"

            if not key:
                continue

            details = group[~group["_is_channel_total"]]
            if details.empty:
                details = group

            est_gross = float(details["Gross spend by channel / platform"].fillna(0).sum())
            est_impressions = float(details["Estimated Impressions"].fillna(0).sum())
            est_clicks = float(details["Estimated link clicks"].fillna(0).sum())
            est_reach = float(details["Estimated Reach"].fillna(0).sum())
            est_ctr_series = details["Estimated CTR"].dropna()
            raw_ctr = coerce_float(est_ctr_series.iloc[0]) if not est_ctr_series.empty else 0.0

            est_cpm_series = details["Net CPM"].dropna()
            est_cpm = coerce_float(est_cpm_series.iloc[0]) if not est_cpm_series.empty else 0.0

            est_freq_series = details["Estimated Frequency"].dropna()
            if not est_freq_series.empty:
                est_freq = coerce_float(est_freq_series.iloc[0])
            else:
                est_freq = safe_div(est_impressions, est_reach)

            est_ctr = raw_ctr * 100 if raw_ctr and raw_ctr <= 1 else raw_ctr
            estimates[key] = {
                "gross_spend": est_gross,
                "impressions": est_impressions,
                "clicks": est_clicks,
                "reach": est_reach,
                "ctr": est_ctr,
                "net_cpm": est_cpm,
                "frequency": est_freq,
            }

        if "_is_channel_total" in df.columns:
            df.drop(columns=["_is_channel_total"], inplace=True, errors="ignore")

        return estimates

    def _build_overall_summary(
        self,
        meta: ChannelSummary,
        pinterest: Optional[ChannelSummary],
        tiktok: Optional[ChannelSummary],
    ) -> Dict[str, float]:
        summaries = [meta]
        if pinterest:
            summaries.append(pinterest)
        if tiktok:
            summaries.append(tiktok)

        total_gross = sum(s.gross_spend for s in summaries)
        total_net = sum(s.net_spend for s in summaries)
        total_impressions = sum(s.impressions for s in summaries)
        total_reach = sum(s.reach for s in summaries)
        total_clicks = sum(s.clicks for s in summaries)

        total_ctr = safe_div(total_clicks, total_impressions) * 100
        total_cpm = safe_div(total_net, total_impressions) * 1000
        total_frequency = safe_div(total_impressions, total_reach)

        brand_revenue = sum(s.brand_revenue for s in summaries)
        brand_units = sum(s.brand_units for s in summaries)
        brand_roas = safe_div(brand_revenue, total_net)
        brand_roi = safe_div(brand_revenue, total_gross)

        fsku_revenue = sum(s.fsku_revenue for s in summaries)
        fsku_units = sum(s.fsku_units for s in summaries)
        fsku_roas = safe_div(fsku_revenue, total_net)
        fsku_roi = safe_div(fsku_revenue, total_gross)

        online_revenue = sum(s.brand_online_revenue for s in summaries)
        instore_revenue = sum(s.brand_instore_revenue for s in summaries)

        return {
            "gross_spend": total_gross,
            "net_spend": total_net,
            "impressions": total_impressions,
            "reach": total_reach,
            "clicks": total_clicks,
            "ctr": total_ctr,
            "net_cpm": total_cpm,
            "frequency": total_frequency,
            "brand_revenue": brand_revenue,
            "brand_units": brand_units,
            "brand_roas": brand_roas,
            "brand_roi": brand_roi,
            "fsku_revenue": fsku_revenue,
            "fsku_units": fsku_units,
            "fsku_roas": fsku_roas,
            "fsku_roi": fsku_roi,
            "brand_online_revenue": online_revenue,
            "brand_instore_revenue": instore_revenue,
        }

    # ------------------------------------------------------------------ #
    # Replacement helpers
    # ------------------------------------------------------------------ #

    def _basic_identity_fields(
        self,
        df: pd.DataFrame,
        media_plan: Optional[pd.DataFrame],
        objectives: Dict[str, str],
    ) -> Dict[str, str]:
        brand_name = df["Brand"].dropna().iloc[0] if not df["Brand"].dropna().empty else "Unknown Brand"
        campaign_name = df["A360_Campaign"].dropna().iloc[0] if not df["A360_Campaign"].dropna().empty else "Campaign"
        campaign_dates = date.today().strftime("%d/%m/%Y")

        flight_dates = campaign_dates
        if media_plan is not None:
            first_flight = media_plan["Flight duration"].dropna()
            if not first_flight.empty:
                flight_dates = str(first_flight.iloc[0]).replace("\n", " ").strip()

        channels_text = self._channels_human_readable()

        display_map = {"meta": "Meta", "pin": "Pinterest", "tik": "TikTok"}
        channels_sorted = [display_map.get(c, c.title()) for c in self.channels_present]
        channel1 = channels_sorted[0] if channels_sorted else ""
        channel2 = channels_sorted[1] if len(channels_sorted) > 1 else ""

        campaign_objective = objectives.get("campaign_objective", "")
        objective_1 = objectives.get("objective_1", "")
        objective_2 = objectives.get("objective_2", "")
        primary_kpis = objectives.get("primary_kpis", "")
        secondary_kpis = objectives.get("secondary_kpis", "")

        return {
            "{brand_name}": brand_name,
            "{campaign_name}": campaign_name,
            "{campaign_objective}": campaign_objective,
            "{campaign_objective_1}": objective_1,
            "{campaign_objective_2}": objective_2,
            "{primary_kpis}": primary_kpis,
            "{secondary_kpis}": secondary_kpis,
            "{date_}": date.today().strftime("%d/%m/%Y"),
            "{campaign_dates}": campaign_dates,
            "{flight_dates}": flight_dates,
            "{channels}": channels_text,
            "{channel1}": channel1,
            "{channel2}": channel2,
        }

    def _overall_metrics(
        self,
        overall: Dict[str, float],
        estimates: Dict[str, Dict[str, float]],
        est_totals: Optional[Dict[str, float]] = None,
    ) -> Dict[str, str]:
        est_totals = est_totals or self._aggregate_estimates(estimates)

        perc_imp = self._calc_delta(overall["impressions"], est_totals["impressions"])
        perc_clicks = self._calc_delta(overall["clicks"], est_totals["clicks"])

        replacements = {
            "{gross_spend}": fmt_currency(overall["gross_spend"]),
            "{impressions}": fmt_int(overall["impressions"]),
            "{reach}": fmt_int(overall["reach"]),
            "{clicks}": fmt_int(overall["clicks"]),
            "{ctr}": fmt_percent(overall["ctr"]),
            "{net_cpm}": fmt_currency(overall["net_cpm"]),
            "{freq_est}": fmt_number(overall["frequency"]),
            "{brand_rev}": fmt_currency(overall["brand_revenue"]),
            "{brand_roas}": fmt_number(overall["brand_roas"]),
            "{brand_roi}": fmt_number(overall["brand_roi"]),
            "{brand_units}": fmt_int(overall["brand_units"]),
            "{fsku_rev}": fmt_currency(overall["fsku_revenue"]),
            "{fsku_roas}": fmt_number(overall["fsku_roas"]),
            "{fsku_roi}": fmt_number(overall["fsku_roi"]),
            "{fsku_units}": fmt_int(overall["fsku_units"]),
            "{gross_est}": fmt_currency(est_totals["gross"]),
            "{imp_est}": fmt_int(est_totals["impressions"]),
            "{click_est}": fmt_int(est_totals["clicks"]),
            "{ctr_est}": fmt_percent(est_totals["ctr"]),
            "{net_cpm_est}": fmt_currency(est_totals["net_cpm"]),
            "{perc_imp}": fmt_percent(perc_imp),
            "{perc_clicks}": fmt_percent(perc_clicks),
        }
        return replacements

    def _channel_placeholders(
        self,
        prefix: str,
        summary: ChannelSummary,
        estimates: Optional[Dict[str, float]],
        audiences: Optional[List[AudiencePerformance]],
        creatives: Optional[List[CreativePerformance]],
    ) -> Dict[str, str]:
        replacements = {
            f"{{{prefix}_gross_spend}}": fmt_currency(summary.gross_spend),
            f"{{{prefix}_net_cpm}}": fmt_currency(summary.net_cpm),
            f"{{{prefix}_impressions}}": fmt_int(summary.impressions),
            f"{{{prefix}_reach}}": fmt_int(summary.reach),
            f"{{{prefix}_clicks}}": fmt_int(summary.clicks),
            f"{{{prefix}_ctr}}": fmt_percent(summary.ctr),
            f"{{{prefix}_freq}}": fmt_number(summary.frequency),
            f"{{{prefix}_brand_revenue}}": fmt_currency(summary.brand_revenue),
            f"{{{prefix}_brand_rev}}": fmt_currency(summary.brand_revenue),
            f"{{{prefix}_brand_roas}}": fmt_number(summary.brand_roas),
            f"{{{prefix}_brand_roi}}": fmt_number(summary.brand_roi),
            f"{{{prefix}_brand_online_revenue}}": fmt_currency(summary.brand_online_revenue),
            f"{{{prefix}_brand_instore_revenue}}": fmt_currency(summary.brand_instore_revenue),
            f"{{{prefix}_brand_online_perc_sales}}": fmt_number(
                safe_div(summary.brand_online_revenue, summary.brand_revenue) * 100 if summary.brand_revenue else 0
            ),
            f"{{{prefix}_brand_instore_perc_sales}}": fmt_number(
                safe_div(summary.brand_instore_revenue, summary.brand_revenue) * 100 if summary.brand_revenue else 0
            ),
            f"{{{prefix}_fsku_revenue}}": fmt_currency(summary.fsku_revenue),
            f"{{{prefix}_fsku_rev}}": fmt_currency(summary.fsku_revenue),
            f"{{{prefix}_fsku_roas}}": fmt_number(summary.fsku_roas),
            f"{{{prefix}_fsku_roi}}": fmt_number(summary.fsku_roi),
            f"{{{prefix}_fsku_online_revenue}}": fmt_currency(summary.fsku_online_revenue),
            f"{{{prefix}_fsku_instore_revenue}}": fmt_currency(summary.fsku_instore_revenue),
            f"{{{prefix}_fsku_units}}": fmt_int(summary.fsku_units),
        }

        if estimates:
            replacements.update(
                {
                    f"{{{prefix}_est_spend}}": fmt_currency(estimates.get("gross_spend", 0)),
                    f"{{{prefix}_est_imp}}": fmt_int(estimates.get("impressions", 0)),
                    f"{{{prefix}_est_reach}}": fmt_int(estimates.get("reach", 0)),
                    f"{{{prefix}_est_clicks}}": fmt_int(estimates.get("clicks", 0)),
                    f"{{{prefix}_est_ctr}}": fmt_percent(estimates.get("ctr", 0)),
                    f"{{{prefix}_est_cpm}}": fmt_currency(estimates.get("net_cpm", 0)),
                    f"{{{prefix}_est_freq}}": fmt_number(estimates.get("frequency", 0)),
                }
            )
        else:
            replacements.update(
                {
                    f"{{{prefix}_est_spend}}": "N/A",
                    f"{{{prefix}_est_imp}}": "N/A",
                    f"{{{prefix}_est_reach}}": "N/A",
                    f"{{{prefix}_est_clicks}}": "N/A",
                    f"{{{prefix}_est_ctr}}": "N/A",
                    f"{{{prefix}_est_cpm}}": "N/A",
                    f"{{{prefix}_est_freq}}": "N/A",
                }
            )

        replacements[f"{{{prefix}_two_audiences_with_strongest_roas_commentary}}"] = self._top_two_audience_blurb(audiences)

        return replacements

    def _empty_channel_placeholders(self, prefix: str) -> Dict[str, str]:
        replacements = {}
        for key in [
            "gross_spend",
            "net_cpm",
            "impressions",
            "reach",
            "clicks",
            "ctr",
            "freq",
            "brand_revenue",
            "brand_rev",
            "brand_roas",
            "brand_roi",
            "brand_online_revenue",
            "brand_instore_revenue",
            "brand_online_perc_sales",
            "brand_instore_perc_sales",
            "fsku_revenue",
            "fsku_rev",
            "fsku_roas",
            "fsku_roi",
            "fsku_online_revenue",
            "fsku_instore_revenue",
            "fsku_units",
            "est_spend",
            "est_imp",
            "est_reach",
            "est_clicks",
            "est_ctr",
            "est_cpm",
            "est_freq",
            "two_audiences_with_strongest_roas_commentary",
        ]:
            replacements[f"{{{prefix}_{key}}}"] = "N/A"
        return replacements

    def _audience_tables(
        self,
        meta_audiences: List[AudiencePerformance],
        pin_audiences: Optional[List[AudiencePerformance]],
        tik_audiences: Optional[List[AudiencePerformance]],
    ) -> Dict[str, str]:
        replacements: Dict[str, str] = {}
        self._fill_audience_block(replacements, "m", meta_audiences)
        self._fill_audience_block(replacements, "p", pin_audiences or [])
        self._fill_audience_block(replacements, "t", tik_audiences or [])
        return replacements

    def _fill_audience_block(self, replacements: Dict[str, str], prefix: str, audiences: List[AudiencePerformance]) -> None:
        for idx in range(1, 13):
            audience = audiences[idx - 1] if idx <= len(audiences) else None
            replacements[f"{{{prefix}_aud_{idx}}}"] = audience.name if audience else "N/A"
            replacements[f"{{{prefix}_net_spend_{idx}}}"] = fmt_currency(audience.net_spend) if audience else "N/A"
            replacements[f"{{{prefix}_netspend_{idx}}}"] = fmt_currency(audience.net_spend) if audience else "N/A"
            replacements[f"{{{prefix}_imp_{idx}}}"] = fmt_int(audience.impressions) if audience else "N/A"
            replacements[f"{{{prefix}_reach_{idx}}}"] = fmt_int(audience.reach) if audience else "N/A"
            replacements[f"{{{prefix}_freq_{idx}}}"] = fmt_number(audience.frequency) if audience else "N/A"
            replacements[f"{{{prefix}_clicks_{idx}}}"] = fmt_int(audience.clicks) if audience else "N/A"
            replacements[f"{{{prefix}_ctr_{idx}}}"] = fmt_percent(audience.ctr) if audience else "N/A"
            replacements[f"{{{prefix}_netcpm_{idx}}}"] = fmt_currency(audience.net_cpm) if audience else "N/A"
            replacements[f"{{{prefix}_revenue_{idx}}}"] = fmt_currency(audience.revenue) if audience else "N/A"
            replacements[f"{{{prefix}_roas_{idx}}}"] = fmt_number(audience.roas) if audience else "N/A"

    def _creative_placeholders(
        self, prefix: str, creatives: Optional[List[CreativePerformance]]
    ) -> Dict[str, str]:
        replacements = {
            f"{{{prefix}_ad_with_strongest_roas_picture}}": "",
            f"{{{prefix}_ad_with_strongest_roas_picture_description}}": "",
            f"{{{prefix}_ad_with_weakest_roas_picture}}": "",
            f"{{{prefix}_ad_with_weakest_roas_picture_description}}": "",
            f"{{{prefix}_ad_with_strongest_roas_and_impressions_commentary}}": "Not applicable",
            f"{{{prefix}_ad_with_weakest_roas_and_impressions_commentary}}": "Not applicable",
        }

        if not creatives:
            return replacements

        channel_suffix = {"meta": "m", "pin": "p", "tik": "t"}.get(prefix, prefix[:1].lower())
        sorted_creatives = sorted(creatives, key=lambda c: c.roas, reverse=True)
        strongest = sorted_creatives[0]
        weakest = sorted_creatives[-1]

        replacements.update(
            {
                f"{{{prefix}_ad_with_strongest_roas_picture_description}}": strongest.name,
                f"{{{prefix}_ad_with_strongest_roas_and_impressions_commentary}}": self._creative_commentary(strongest, True),
                f"{{{prefix}_ad_with_weakest_roas_picture_description}}": weakest.name,
                f"{{{prefix}_ad_with_weakest_roas_and_impressions_commentary}}": self._creative_commentary(weakest, False),
            }
        )

        replacements.update(
            self._creative_metric_block("adwithstrongestroas", channel_suffix, strongest)
        )
        replacements.update(
            self._creative_metric_block("adwithweakestroas", channel_suffix, weakest)
        )

        return replacements

    def _creative_commentary(self, creative: CreativePerformance, is_strongest: bool) -> str:
        adjective = "strongest" if is_strongest else "weakest"
        roas_text = f"ROAS {fmt_number(creative.roas)}" if creative.roas else "ROAS unavailable"
        return (
            f"{creative.name} delivered {fmt_int(creative.impressions)} impressions at "
            f"{fmt_percent(creative.ctr)} CTR ({roas_text})."
        )

    def _creative_metric_block(
        self,
        base_placeholder: str,
        suffix: str,
        creative: CreativePerformance,
    ) -> Dict[str, str]:
        return {
            f"{{{base_placeholder}_imp_{suffix}}}": fmt_int(creative.impressions),
            f"{{{base_placeholder}_reach_{suffix}}}": fmt_int(creative.reach),
            f"{{{base_placeholder}_spend_{suffix}}}": fmt_currency(creative.net_spend),
            f"{{{base_placeholder}_freq_{suffix}}}": fmt_number(creative.frequency),
            f"{{{base_placeholder}_roas_{suffix}}}": fmt_number(creative.roas),
            f"{{{base_placeholder}_cpm_{suffix}}}": fmt_currency(creative.net_cpm),
            f"{{{base_placeholder}_clicks_{suffix}}}": fmt_int(creative.clicks),
            f"{{{base_placeholder}_ctr_{suffix}}}": fmt_percent(creative.ctr),
        }

    def _top_two_audience_blurb(self, audiences: Optional[List[AudiencePerformance]]) -> str:
        if not audiences:
            return "Not applicable"
        top = sorted(audiences, key=lambda a: a.roas, reverse=True)[:2]
        descriptions = [
            f"{aud.name} (ROAS {fmt_number(aud.roas)})"
            for aud in top
            if aud.roas
        ]
        return ", ".join(descriptions) if descriptions else "Not applicable"

    def _calc_delta(self, actual: float, expected: float) -> float:
        if expected == 0:
            return 0.0
        return safe_div(actual - expected, expected) * 100

    def _aggregate_estimates(self, estimates: Dict[str, Dict[str, float]]) -> Dict[str, float]:
        if not estimates:
            return {
                "gross": 0.0,
                "impressions": 0.0,
                "clicks": 0.0,
                "reach": 0.0,
                "ctr": 0.0,
                "net_cpm": 0.0,
                "frequency": 0.0,
            }

        totals = {
            "gross": 0.0,
            "impressions": 0.0,
            "clicks": 0.0,
            "reach": 0.0,
            "ctr_values": [],
            "net_cpm_values": [],
            "frequency_values": [],
        }

        for data in estimates.values():
            totals["gross"] += data.get("gross_spend", 0) or 0.0
            totals["impressions"] += data.get("impressions", 0) or 0.0
            totals["clicks"] += data.get("clicks", 0) or 0.0
            totals["reach"] += data.get("reach", 0) or 0.0

            if data.get("ctr") is not None:
                totals["ctr_values"].append(data["ctr"])
            if data.get("net_cpm") is not None:
                totals["net_cpm_values"].append(data["net_cpm"])
            if data.get("frequency") is not None:
                totals["frequency_values"].append(data["frequency"])

        return {
            "gross": totals["gross"],
            "impressions": totals["impressions"],
            "clicks": totals["clicks"],
            "reach": totals["reach"],
            "ctr": float(np.mean(totals["ctr_values"])) if totals["ctr_values"] else 0.0,
            "net_cpm": float(np.mean(totals["net_cpm_values"])) if totals["net_cpm_values"] else 0.0,
            "frequency": float(np.mean(totals["frequency_values"])) if totals["frequency_values"] else 0.0,
        }

    def _generate_summary_charts(self, overall: Dict[str, float], estimate_totals: Dict[str, float]) -> None:
        try:
            generate_summary_charts(
                impressions_actual=overall.get("impressions", 0.0),
                impressions_estimate=estimate_totals.get("impressions", 0.0),
                clicks_actual=overall.get("clicks", 0.0),
                clicks_estimate=estimate_totals.get("clicks", 0.0),
                output_dir=os.getcwd(),
            )
        except Exception as exc:  # pragma: no cover - best effort
            print(f"⚠️ Unable to generate summary charts: {exc}")

    def _channels_human_readable(self) -> str:
        names = []
        mapping = {"meta": "Meta", "pin": "Pinterest", "tik": "TikTok"}
        for channel in self.channels_present:
            names.append(mapping.get(channel, channel.title()))
        if len(names) == 1:
            return names[0]
        if len(names) == 2:
            return " and ".join(names)
        return ", ".join(names[:-1]) + f", and {names[-1]}"

    def _online_instore_metrics(
        self,
        meta: ChannelSummary,
        pinterest: Optional[ChannelSummary],
    ) -> Tuple[Dict[str, str], Dict[str, float]]:
        online_revenue = meta.brand_online_revenue + (pinterest.brand_online_revenue if pinterest else 0)
        instore_revenue = meta.brand_instore_revenue + (pinterest.brand_instore_revenue if pinterest else 0)
        online_units = meta.brand_online_units
        instore_units = meta.brand_instore_units
        if pinterest:
            online_units += pinterest.brand_online_units
            instore_units += pinterest.brand_instore_units

        online_aov = safe_div(online_revenue, online_units)
        instore_aov = safe_div(instore_revenue, instore_units)
        total_revenue = online_revenue + instore_revenue
        online_share = safe_div(online_revenue, total_revenue) * 100 if total_revenue else 0
        instore_share = safe_div(instore_revenue, total_revenue) * 100 if total_revenue else 0

        replacements = {
            "{online_revenue}": fmt_currency(online_revenue),
            "{instore_aov}": fmt_currency(instore_aov),
            "{online_aov}": fmt_currency(online_aov),
            "{online_units}": fmt_int(online_units),
            "{instore_units}": fmt_int(instore_units),
        }

        metrics = {
            "online_revenue": online_revenue,
            "instore_revenue": instore_revenue,
            "total_revenue": total_revenue,
            "online_units": online_units,
            "instore_units": instore_units,
            "online_aov": online_aov,
            "instore_aov": instore_aov,
            "online_share_percent": online_share,
            "instore_share_percent": instore_share,
        }

        return replacements, metrics

    # ------------------------------------------------------------------ #
    # Commentaries via GPT
    # ------------------------------------------------------------------ #

    def _generate_commentaries(
        self,
        overall: Dict[str, float],
        estimates: Dict[str, Dict[str, float]],
        channel_context: Dict[str, Dict[str, object]],
        online_instore_metrics: Dict[str, float],
        objectives: Dict[str, str],
    ) -> Dict[str, str]:
        if client is None:
            return {key: "OpenAI API key not configured." for key in self._commentary_keys()}

        payload = {
            "overall": overall,
            "estimates": estimates,
            "channels_present": self.channels_present,
            "channels": {},
            "online_instore": online_instore_metrics,
            "objectives": objectives,
        }

        for channel, ctx in channel_context.items():
            audiences: List[AudiencePerformance] = ctx.get("audiences", [])
            creatives: List[CreativePerformance] = ctx.get("creatives", [])
            summary: Optional[ChannelSummary] = ctx.get("summary")

            audiences_sorted = sorted(
                audiences,
                key=lambda a: (
                    self._sanitize_value(a.roas),
                    self._sanitize_value(a.revenue),
                ),
                reverse=True,
            )

            top_creatives = self._creative_snapshot(creatives, limit=3, best=True)
            best_creative_entry = self._creative_snapshot(creatives, limit=1, best=True)
            worst_creative_entry = self._creative_snapshot(creatives, limit=1, best=False)

            payload["channels"][channel] = {
                "summary": self._summary_snapshot(summary),
                "audience_count": len(audiences),
                "top_two_audiences": self._audience_snapshot(audiences_sorted, limit=2),
                "audiences": self._audience_snapshot(audiences, limit=12),
                "top_creatives": top_creatives,
                "best_creative": best_creative_entry[0] if best_creative_entry else None,
                "worst_creative": worst_creative_entry[0] if worst_creative_entry else None,
                "creative_count": len(creatives),
            }

        prompt = self._build_commentary_prompt(payload)
        response_text = self.gpt4_completion(prompt)

        try:
            commentary_json = json.loads(response_text)
        except json.JSONDecodeError:
            commentary_json = {key: response_text for key in self._commentary_keys()}

        replacements = {}

        keys = self._commentary_keys()
        missing_channels = {
            "meta": "meta" not in self.channels_present,
            "pin": "pin" not in self.channels_present,
            "tik": "tik" not in self.channels_present,
        }

        for key in keys:
            value = commentary_json.get(key, "Not applicable")
            if not isinstance(value, str):
                value = json.dumps(value)

            for channel, should_blank in missing_channels.items():
                if should_blank and key in self.CHANNEL_COMMENTARY_KEYS[channel]:
                    value = "Not applicable"
                    break

            replacements[f"{{{key}}}"] = value.strip()

        return replacements

    def _build_commentary_prompt(self, payload: Dict[str, object]) -> str:
        keys = self._commentary_keys()
        channels = ", ".join(self.channels_present)

        if self.prompt_template and Path(self.prompt_template).exists():
            template_text = Path(self.prompt_template).read_text()
            prompt = (
                template_text.replace("<<CHANNELS>>", channels)
                .replace("<<PAYLOAD>>", json.dumps(payload, default=str))
                .replace("<<KEY_LIST>>", "\n".join(f"- {{{key}}}" for key in keys))
                .replace("<<KEYS_COMMA>>", ", ".join(f"{{{key}}}" for key in keys))
            )
            return prompt

        return (
            "You are an experienced marketing analyst. Craft concise, natural-language insights "
            "based on the supplied campaign data payload. Output must be a single JSON object "
            "where every value is a short paragraph (one to two sentences) with no nested JSON.\n"
            f"Channels present: {channels or 'None'}\n"
            f"Data payload:\n{json.dumps(payload, default=str)}\n\n"
            "Return strings for the following keys:\n"
            + "\n".join(f"- {{{key}}}" for key in keys)
            + "\n\nGuidelines:\n"
            "- Use plain sentences (no bullet lists).\n"
            "- Reference specific metrics when available; otherwise describe the trend plainly.\n"
            "- If a key relates to a channel that is absent, respond with 'Not applicable'.\n"
            "- For ad-level commentaries mention the creative name and note why it leads or lags.\n"
        )

    def _audience_snapshot(
        self,
        audiences: List[AudiencePerformance],
        limit: Optional[int] = 5,
    ) -> List[Dict[str, object]]:
        selected = audiences if limit is None else audiences[:limit]
        snapshot: List[Dict[str, object]] = []
        for aud in selected:
            snapshot.append(
                {
                    "name": aud.name,
                    "net_spend": self._sanitize_value(aud.net_spend),
                    "impressions": self._sanitize_value(aud.impressions),
                    "reach": self._sanitize_value(aud.reach),
                    "frequency": self._sanitize_value(aud.frequency),
                    "clicks": self._sanitize_value(aud.clicks),
                    "ctr": self._sanitize_value(aud.ctr),
                    "net_cpm": self._sanitize_value(aud.net_cpm),
                    "revenue": self._sanitize_value(aud.revenue),
                    "roas": self._sanitize_value(aud.roas),
                    "roi": self._sanitize_value(aud.roi),
                }
            )
        return snapshot

    def _creative_snapshot(
        self,
        creatives: List[CreativePerformance],
        limit: Optional[int] = 3,
        best: bool = True,
    ) -> List[Dict[str, object]]:
        if not creatives:
            return []

        sorted_creatives = sorted(
            creatives,
            key=lambda c: (
                self._sanitize_value(c.roas),
                self._sanitize_value(c.impressions),
            ),
            reverse=best,
        )

        selected = sorted_creatives if limit is None else sorted_creatives[:limit]

        snapshot: List[Dict[str, object]] = []
        for creative in selected:
            snapshot.append(
                {
                    "name": creative.name,
                    "net_spend": self._sanitize_value(creative.net_spend),
                    "impressions": self._sanitize_value(creative.impressions),
                    "reach": self._sanitize_value(creative.reach),
                    "frequency": self._sanitize_value(creative.frequency),
                    "clicks": self._sanitize_value(creative.clicks),
                    "ctr": self._sanitize_value(creative.ctr),
                    "net_cpm": self._sanitize_value(creative.net_cpm),
                    "revenue": self._sanitize_value(creative.revenue),
                    "roas": self._sanitize_value(creative.roas),
                }
            )
        return snapshot

    def _summary_snapshot(self, summary: Optional[ChannelSummary]) -> Dict[str, object]:
        if summary is None:
            return {}
        data = asdict(summary)
        result: Dict[str, object] = {}
        for key, value in data.items():
            if isinstance(value, (int, float, np.floating)):
                result[key] = self._sanitize_value(float(value))
            else:
                result[key] = value
        return result

    def _extract_objectives(self, meta_df: pd.DataFrame) -> Dict[str, str]:
        objectives: List[str] = []
        if "Campaign_Object" in meta_df.columns:
            raw = meta_df["Campaign_Object"].dropna().unique()
            objectives = [str(obj).strip() for obj in raw if str(obj).strip()]
        while len(objectives) < 2:
            objectives.append("Not specified")
        objective_1 = objectives[0] if objectives else ""
        objective_2 = objectives[1] if len(objectives) > 1 else ""

        result = {
            "objective_1": objective_1,
            "objective_2": objective_2,
            "campaign_objective": objective_1,
            "primary_kpis": "",
            "secondary_kpis": "",
        }

        if self.manual_campaign_objective:
            result["campaign_objective"] = self.manual_campaign_objective
            result["objective_1"] = self.manual_campaign_objective

        if self.manual_primary_kpis:
            result["primary_kpis"] = self.manual_primary_kpis
        if self.manual_secondary_kpis:
            result["secondary_kpis"] = self.manual_secondary_kpis

        return result

    @staticmethod
    def _sanitize_value(value: object) -> object:
        if value is None:
            return 0.0
        if isinstance(value, np.generic):
            value = float(value)
        if isinstance(value, (int, float)):
            value = float(value)
            if math.isnan(value):
                return 0.0
            return value
        return value

    def _commentary_keys(self) -> List[str]:
        return list(self.COMMENTARY_KEYS_IN_ORDER)

    def gpt4_completion(self, prompt: str, engine: str = "gpt-4o-mini", temp: float = 0.3, tokens: int = 2000) -> str:
        max_retry = 3
        retry = 0
        while retry < max_retry:
            try:
                response = client.responses.create(
                    model=engine,
                    input=prompt,
                    temperature=temp,
                    max_output_tokens=tokens,
                )
                return response.output_text
            except Exception as exc:  # pragma: no cover - external API
                retry += 1
                if retry >= max_retry:
                    return f"Unable to generate commentary: {exc}"

    # ------------------------------------------------------------------ #
    # Placeholder completion
    # ------------------------------------------------------------------ #

    def _fill_missing_placeholders(self, replacements: Dict[str, str]) -> Dict[str, str]:
        try:
            from pptx import Presentation
        except ImportError:  # pragma: no cover
            return replacements

        template_path = Path(self.template_path)
        if not template_path.exists():
            return replacements

        prs = Presentation(str(template_path))
        all_placeholders: set[str] = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                texts = []
                if shape.has_text_frame:
                    texts.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
                for text in texts:
                    all_placeholders.update(self._extract_placeholders(text))

        for placeholder in all_placeholders:
            if placeholder not in replacements:
                replacements[placeholder] = "N/A"

        return replacements

    @staticmethod
    def _extract_placeholders(text: str) -> Iterable[str]:
        start = 0
        results = []
        while True:
            open_idx = text.find("{", start)
            if open_idx == -1:
                break
            close_idx = text.find("}", open_idx)
            if close_idx == -1:
                break
            results.append(text[open_idx : close_idx + 1])
            start = close_idx + 1
        return results


# Convenience alias for existing imports
DataExtractor = ConsolidatedDataExtractor
